from __future__ import annotations

import os
import json
import sys
import asyncio
import subprocess
import tempfile
import zipfile
import math
import random
import wave
from collections import deque
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree

import fitz
import cv2
from docx import Document
from pdf2docx import Converter as PdfToDocxConverter
from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError
from pdf2image import convert_from_path
from PIL import Image, ImageDraw, ImageFilter, ImageOps
from pptx import Presentation
from pptx.util import Emu, Pt
from pypdf import PdfReader, PdfWriter
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas

pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
EMU_PER_POINT = 12700


def main() -> int:
    if len(sys.argv) < 2:
        raise SystemExit("missing conversion command")

    command = sys.argv[1]

    if command == "images_to_pdf":
      return images_to_pdf(sys.argv[2], sys.argv[3:])
    if command == "pdf_to_images":
      return pdf_to_images(sys.argv[2], sys.argv[3])
    if command == "pdf_to_word":
      return pdf_to_word(sys.argv[2], sys.argv[3], sys.argv[4], sys.argv[5], sys.argv[6])
    if command == "pdf_to_pptx":
      return pdf_to_pptx(sys.argv[2], sys.argv[3], sys.argv[4] if len(sys.argv) > 4 else "")
    if command == "ocr_text_extract":
      return ocr_text_extract(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "batch_file_rename":
      return batch_file_rename(sys.argv[2], sys.argv[3], sys.argv[4:])
    if command == "payment_code_merge":
      return payment_code_merge(sys.argv[2], sys.argv[3], sys.argv[4:])
    if command == "delete_pages_pdf":
      return delete_pages_pdf(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "reorder_pages_pdf":
      return reorder_pages_pdf(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "protect_unlock_pdf":
      return protect_unlock_pdf(sys.argv[2], sys.argv[3], sys.argv[4], sys.argv[5])
    if command == "watermark_pdf":
      return watermark_pdf(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "add_page_numbers_pdf":
      return add_page_numbers_pdf(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "sign_stamp_pdf":
      return sign_stamp_pdf(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "batch_sign_stamp_pdf":
      return batch_sign_stamp_pdf(sys.argv[2], sys.argv[3], sys.argv[4:])
    if command == "rotate_pdf":
      return rotate_pdf(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "merge_pdf":
      return merge_pdf(sys.argv[2], sys.argv[3:])
    if command == "compress_pdf":
      return compress_pdf(sys.argv[2], sys.argv[3], sys.argv[4], sys.argv[5])
    if command == "pdf_extract_pages":
      return pdf_extract_pages(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "split_pdf":
      return split_pdf(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "docx_to_pdf_fallback":
      return docx_to_pdf_fallback(sys.argv[2], sys.argv[3])
    if command == "zip_files":
      return zip_files(sys.argv[2], sys.argv[3:])
    if command == "image_compress_batch":
      return image_compress_batch(sys.argv[2], sys.argv[3], sys.argv[4:])
    if command == "image_resize_exact":
      return image_resize_exact(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_resize_scale":
      return image_resize_scale(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_crop_free":
      return image_crop_free(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_crop_ratio":
      return image_crop_ratio(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_crop_ratio_batch":
      return image_crop_ratio_batch(sys.argv[2], sys.argv[3], sys.argv[4:])
    if command == "image_split_grid":
      return image_split_grid(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_concat_long":
      return image_concat_long(sys.argv[2], sys.argv[3], sys.argv[4:])
    if command == "image_collage":
      return image_collage(sys.argv[2], sys.argv[3], sys.argv[4:])
    if command == "image_fill_background":
      return image_fill_background(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_dark_mode_background":
      return image_fill_background(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_watermark_tile":
      return image_watermark_tile(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_grayscale":
      return image_grayscale(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_invert":
      return image_invert(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_printmaking":
      return image_printmaking(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_emboss":
      return image_emboss(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_remove_solid_bg":
      return image_remove_solid_bg(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_smart_bg_remove":
      return image_smart_bg_remove(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "qr_generate":
      return qr_generate(sys.argv[2], sys.argv[3])
    if command == "qr_generate_batch":
      return qr_generate_batch(sys.argv[2], sys.argv[3])
    if command == "qr_decode":
      return qr_decode(sys.argv[2], sys.argv[3])
    if command == "favicon_generate":
      return favicon_generate(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "app_icon_generate":
      return app_icon_generate(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "chrome_icon_generate":
      return chrome_icon_generate(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_add_padding":
      return image_add_padding(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_pixelate":
      return image_pixelate(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_increase_size":
      return image_increase_size(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_clear_content":
      return image_clear_content(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_heic_convert":
      return image_heic_convert(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_format_convert":
      return image_format_convert_dispatch(sys.argv[2], sys.argv[3:])
    if command == "excel_extract_images":
      return extract_images_from_office_zip(sys.argv[2], sys.argv[3], "xl/media/")
    if command == "ppt_extract_images":
      return extract_images_from_office_zip(sys.argv[2], sys.argv[3], "ppt/media/")
    if command == "image_modify_dpi":
      return image_modify_dpi(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "gif_split":
      return gif_split(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "gif_merge":
      return gif_merge(sys.argv[2], sys.argv[3], sys.argv[4:])
    if command == "png_alpha_invert":
      return png_alpha_invert(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_round_corner":
      return image_round_corner(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "image_tile_fill":
      return image_tile_fill(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "id_photo_resize":
      return id_photo_resize(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "exam_id_photo_process":
      return id_photo_resize(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "id_photo_crop":
      return id_photo_crop(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "id_photo_bg_swap":
      return id_photo_bg_swap(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "anti_ocr_image":
      return anti_ocr_image(sys.argv[2], sys.argv[3], sys.argv[4])
    if command == "text_to_speech":
      return text_to_speech(
          sys.argv[2],
          sys.argv[3],
          sys.argv[4],
          sys.argv[5],
          sys.argv[6] if len(sys.argv) > 6 else "",
      )
    if command == "audio_to_text":
      return audio_to_text(
          sys.argv[2],
          sys.argv[3],
          sys.argv[4] if len(sys.argv) > 4 else "auto",
          sys.argv[5] if len(sys.argv) > 5 else "",
      )

    raise SystemExit(f"unsupported conversion command: {command}")


def images_to_pdf(output_path: str, input_paths: list[str]) -> int:
    if not input_paths:
        raise SystemExit("images_to_pdf requires at least one image")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    first_image = Image.open(input_paths[0])
    page_width, page_height = first_image.size
    pdf = canvas.Canvas(str(output), pagesize=(page_width, page_height))

    for index, image_path in enumerate(input_paths):
        image = Image.open(image_path)
        width, height = image.size
        if index > 0:
            pdf.setPageSize((width, height))
        pdf.drawImage(ImageReader(image), 0, 0, width=width, height=height)
        pdf.showPage()

    pdf.save()
    return 0


def pdf_to_images(output_prefix: str, pdf_path: str) -> int:
    poppler_bin_dir = os.environ.get("POPPLER_BIN_DIR") or None
    try:
        prefix = Path(output_prefix)
        prefix.parent.mkdir(parents=True, exist_ok=True)

        image_paths = convert_from_path(
            pdf_path,
            poppler_path=poppler_bin_dir,
            dpi=144,
            fmt="png",
            output_folder=str(prefix.parent),
            output_file=prefix.name,
            paths_only=True,
            thread_count=1,
        )
    except PDFInfoNotInstalledError as exc:
        raise SystemExit(
            "pdf_to_images requires poppler (pdfinfo/pdftoppm) to be installed and reachable"
        ) from exc
    except PDFPageCountError as exc:
        raise SystemExit("Unable to read PDF page count for pdf_to_images") from exc

    if not image_paths:
        raise SystemExit("pdf_to_images did not render any pages")

    return 0


def pdf_to_word(output_path: str, input_path: str, conversion_mode: str, ocrmypdf_bin: str, ocr_language: str) -> int:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    source_pdf_path = input_path

    with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
        if conversion_mode == "ocr":
            if not ocrmypdf_bin:
                raise SystemExit("OCR mode requires OCRmyPDF to be configured")

            ocr_output_path = str(Path(temp_directory) / "ocr-output.pdf")
            try:
                subprocess.run(
                    [
                        ocrmypdf_bin,
                        "--force-ocr",
                        "--language",
                        ocr_language,
                        input_path,
                        ocr_output_path,
                    ],
                    check=True,
                )
            except FileNotFoundError as exc:
                raise SystemExit("OCRmyPDF is not installed or not reachable") from exc
            except subprocess.CalledProcessError as exc:
                raise SystemExit(f"OCRmyPDF failed with exit code {exc.returncode}") from exc

            source_pdf_path = ocr_output_path

        converter = PdfToDocxConverter(source_pdf_path)
        try:
            converter.convert(str(output))
        finally:
            converter.close()

    return 0


def pdf_to_pptx(output_path: str, input_path: str, ocrmypdf_bin: str) -> int:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
        source_pdf_path = input_path
        if ocrmypdf_bin and should_run_ocr_for_pdf_to_pptx(input_path):
            ocr_output_path = str(Path(temp_directory) / "ocr-output.pdf")
            try:
                subprocess.run(
                    [
                        ocrmypdf_bin,
                        "--force-ocr",
                        "--language",
                        "chi_sim+eng",
                        input_path,
                        ocr_output_path,
                    ],
                    check=True,
                )
            except FileNotFoundError as exc:
                raise SystemExit("OCRmyPDF is not installed or not reachable") from exc
            except subprocess.CalledProcessError as exc:
                raise SystemExit(f"OCRmyPDF failed with exit code {exc.returncode}") from exc

            source_pdf_path = ocr_output_path

        build_presentation_from_pdf(source_pdf_path, output)

    return 0


def ocr_text_extract(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    tesseract_bin = str(options.get("tesseractBin") or "").strip()
    tesseract_script_path = str(options.get("tesseractScriptPath") or "").strip()
    ocr_language = str(options.get("ocrLanguage") or "chi_sim+eng").strip() or "chi_sim+eng"
    if not tesseract_bin:
        raise SystemExit("ocr_text_extract requires tesseractBin")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
        temp_root = Path(temp_directory)
        text_base = temp_root / "ocr-output"
        try:
            command = [
                tesseract_bin,
                input_path,
                str(text_base),
                "-l",
                ocr_language,
            ]
            if tesseract_script_path:
                command = [
                    tesseract_bin,
                    tesseract_script_path,
                    input_path,
                    str(text_base),
                    "-l",
                    ocr_language,
                ]
            subprocess.run(
                command,
                check=True,
            )
        except FileNotFoundError as exc:
            raise SystemExit("Tesseract is not installed or not reachable") from exc
        except subprocess.CalledProcessError as exc:
            raise SystemExit(f"Tesseract failed with exit code {exc.returncode}") from exc

        recognized_path = text_base.with_suffix(".txt")
        if not recognized_path.exists():
            raise SystemExit("Tesseract did not produce a txt output")

        output.write_text(recognized_path.read_text(encoding="utf8", errors="ignore"), encoding="utf8")

    return 0


def batch_file_rename(zip_path: str, options_json: str, input_paths: list[str]) -> int:
    if not input_paths:
        raise SystemExit("batch_file_rename requires at least one file")

    options = json.loads(options_json or "{}")
    template = str(options.get("template") or "资料-{n}-{name}")
    start_number = max(1, int(options.get("startNumber") or 1))
    number_width = max(1, min(8, int(options.get("numberWidth") or 2)))

    output = Path(zip_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    used_names: set[str] = set()

    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for index, input_path in enumerate(input_paths):
            source = Path(input_path)
            rendered_number = str(start_number + index).zfill(number_width)
            target_stem = template.replace("{n}", rendered_number).replace("{name}", source.stem).strip()
            safe_stem = sanitize_archive_stem(target_stem or f"file-{rendered_number}")
            arcname = ensure_unique_archive_name(f"{safe_stem}{source.suffix.lower()}", used_names)
            archive.write(source, arcname=arcname)

    return 0


def payment_code_merge(output_path: str, options_json: str, input_paths: list[str]) -> int:
    if len(input_paths) < 2:
        raise SystemExit("payment_code_merge requires at least two images")

    options = json.loads(options_json or "{}")
    layout = "horizontal" if options.get("layout") == "horizontal" else "vertical"
    title = str(options.get("mainTitle") or "收款码").strip()
    gap = 28
    card_padding = 32
    title_height = 84 if title else 32

    images = [Image.open(path).convert("RGBA") for path in input_paths[:4]]
    try:
        target_width = max(image.width for image in images)
        target_height = max(image.height for image in images)
        prepared = [ImageOps.contain(image, (target_width, target_height), Image.Resampling.LANCZOS) for image in images]

        if layout == "horizontal":
            canvas_width = (target_width * len(prepared)) + (gap * (len(prepared) - 1)) + (card_padding * 2)
            canvas_height = target_height + title_height + (card_padding * 2)
        else:
            canvas_width = target_width + (card_padding * 2)
            canvas_height = title_height + (target_height * len(prepared)) + (gap * (len(prepared) - 1)) + (card_padding * 2)

        canvas = Image.new("RGBA", (canvas_width, canvas_height), (255, 255, 255, 255))
        draw = ImageDraw.Draw(canvas)
        if title:
            draw.text((card_padding, 24), title, fill=(17, 24, 39, 255))

        cursor_x = card_padding
        cursor_y = card_padding + title_height
        for image in prepared:
            if layout == "horizontal":
                canvas.alpha_composite(image, (cursor_x, cursor_y + ((target_height - image.height) // 2)))
                cursor_x += target_width + gap
            else:
                canvas.alpha_composite(image, (cursor_x + ((target_width - image.width) // 2), cursor_y))
                cursor_y += target_height + gap

        save_image_file(canvas, Path(output_path), force_format="PNG")
    finally:
        for image in images:
            image.close()

    return 0


def delete_pages_pdf(output_path: str, input_path: str, selection_json: str) -> int:
    payload = json.loads(selection_json)
    deleted_pages = set(payload.get("orderedPages") or [])
    reader = PdfReader(input_path)
    writer = PdfWriter()

    for index, page in enumerate(reader.pages, start=1):
        if index in deleted_pages:
            continue
        writer.add_page(page)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def should_run_ocr_for_pdf_to_pptx(input_path: str) -> bool:
    with fitz.open(input_path) as document:
        if document.page_count == 0:
            return False

        for page in document:
            if page.get_text("text", sort=True).strip():
                return False

    return True


def build_presentation_from_pdf(input_path: str, output_path: Path) -> None:
    with fitz.open(input_path) as document:
        if document.page_count == 0:
            raise SystemExit("pdf_to_pptx requires at least one PDF page")

        first_page = document[0]
        presentation = Presentation()
        presentation.slide_width = Emu(int(first_page.rect.width * EMU_PER_POINT))
        presentation.slide_height = Emu(int(first_page.rect.height * EMU_PER_POINT))
        blank_layout = presentation.slide_layouts[6]

        for page in document:
            slide = presentation.slides.add_slide(blank_layout)
            added_content = add_page_content_to_slide(
                page,
                slide,
                int(presentation.slide_width),
                int(presentation.slide_height),
            )
            if not added_content:
                add_page_snapshot_to_slide(
                    page,
                    slide,
                    int(presentation.slide_width),
                    int(presentation.slide_height),
                )

        presentation.save(str(output_path))


def add_page_content_to_slide(page, slide, slide_width_emu: int, slide_height_emu: int) -> bool:
    page_dict = page.get_text("dict", sort=True)
    added_content = False

    for block in page_dict.get("blocks", []):
        block_type = block.get("type")
        if block_type == 0:
            added_content = add_text_block_to_slide(
                block,
                slide,
                float(page.rect.width),
                float(page.rect.height),
                slide_width_emu,
                slide_height_emu,
            ) or added_content
        elif block_type == 1:
            added_content = add_image_block_to_slide(
                block,
                slide,
                float(page.rect.width),
                float(page.rect.height),
                slide_width_emu,
                slide_height_emu,
            ) or added_content

    return added_content


def add_text_block_to_slide(
    block: dict,
    slide,
    page_width: float,
    page_height: float,
    slide_width_emu: int,
    slide_height_emu: int,
) -> bool:
    lines: list[str] = []
    max_font_size = 12.0

    for line in block.get("lines", []):
        spans = line.get("spans", [])
        line_text = "".join(span.get("text", "") for span in spans).strip()
        if line_text:
            lines.append(line_text)
        for span in spans:
            max_font_size = max(max_font_size, float(span.get("size") or 12))

    if not lines:
        return False

    left, top, width, height = scale_bbox_to_slide(
        block.get("bbox") or (0, 0, 1, 1),
        page_width,
        page_height,
        slide_width_emu,
        slide_height_emu,
    )
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.fill.background()
    textbox.line.fill.background()
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    for index, line_text in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line_text
        paragraph.font.size = Pt(max(10, min(28, max_font_size)))

    return True


def add_image_block_to_slide(
    block: dict,
    slide,
    page_width: float,
    page_height: float,
    slide_width_emu: int,
    slide_height_emu: int,
) -> bool:
    image_bytes = block.get("image")
    if not image_bytes:
        return False

    left, top, width, height = scale_bbox_to_slide(
        block.get("bbox") or (0, 0, 1, 1),
        page_width,
        page_height,
        slide_width_emu,
        slide_height_emu,
    )
    slide.shapes.add_picture(BytesIO(bytes(image_bytes)), left, top, width=width, height=height)
    return True


def add_page_snapshot_to_slide(page, slide, slide_width_emu: int, slide_height_emu: int) -> None:
    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    slide.shapes.add_picture(
        BytesIO(pixmap.tobytes("png")),
        Emu(0),
        Emu(0),
        width=Emu(slide_width_emu),
        height=Emu(slide_height_emu),
    )


def scale_bbox_to_slide(
    bbox, page_width: float, page_height: float, slide_width_emu: int, slide_height_emu: int
) -> tuple[Emu, Emu, Emu, Emu]:
    x0, y0, x1, y1 = bbox
    scale_x = slide_width_emu / max(page_width, 1.0)
    scale_y = slide_height_emu / max(page_height, 1.0)

    left = max(0, int(x0 * scale_x))
    top = max(0, int(y0 * scale_y))
    width = max(1, int((x1 - x0) * scale_x))
    height = max(1, int((y1 - y0) * scale_y))
    return Emu(left), Emu(top), Emu(width), Emu(height)


def reorder_pages_pdf(output_path: str, input_path: str, selection_json: str) -> int:
    payload = json.loads(selection_json)
    ordered_pages = payload.get("orderedPages") or []
    if not ordered_pages:
        raise SystemExit("reorder_pages_pdf requires ordered pages")

    reader = PdfReader(input_path)
    total_pages = len(reader.pages)
    writer = PdfWriter()

    for page_number in ordered_pages:
        ensure_page_in_bounds(page_number, total_pages)
        writer.add_page(reader.pages[page_number - 1])

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def protect_unlock_pdf(output_path: str, input_path: str, mode: str, password: str) -> int:
    reader = PdfReader(input_path)
    writer = PdfWriter()

    if mode == "unlock":
      if not reader.is_encrypted:
        raise SystemExit("The PDF is not password-protected")
      decrypt_result = reader.decrypt(password)
      if decrypt_result == 0:
        raise SystemExit("Incorrect PDF password")

    for page in reader.pages:
        writer.add_page(page)

    if mode == "protect":
        writer.encrypt(password)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def docx_to_pdf_fallback(input_path: str, output_path: str) -> int:
    document = Document(input_path)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    pdf = canvas.Canvas(str(output))
    page_width, page_height = pdf._pagesize
    cursor_y = page_height - 56
    line_height = 20

    title = extract_docx_title(input_path)
    if title:
        pdf.setFont("Helvetica-Bold", 16)
        pdf.drawString(56, cursor_y, title[:90])
        cursor_y -= 32

    pdf.setFont("Helvetica", 11)

    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    if not paragraphs:
        paragraphs = ["(Empty document)"]

    for paragraph in paragraphs:
        for line in wrap_text(paragraph, 90):
            if cursor_y < 56:
                pdf.showPage()
                pdf.setFont("Helvetica", 11)
                cursor_y = page_height - 56
            pdf.drawString(56, cursor_y, line)
            cursor_y -= line_height

        cursor_y -= 4

    pdf.save()
    return 0


def merge_pdf(output_path: str, input_paths: list[str]) -> int:
    if len(input_paths) < 2:
        raise SystemExit("merge_pdf requires at least two PDF files")

    writer = PdfWriter()
    for input_path in input_paths:
        reader = PdfReader(input_path)
        for page in reader.pages:
            writer.add_page(page)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def compress_pdf(output_path: str, input_path: str, compression_level: str, ghostscript_bin: str) -> int:
    if compression_level not in {"standard", "strong"}:
        raise SystemExit("compress_pdf requires compression_level to be standard or strong")

    pdf_setting = "/ebook" if compression_level == "standard" else "/screen"
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    command = [
        ghostscript_bin,
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        f"-dPDFSETTINGS={pdf_setting}",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
        f"-sOutputFile={output}",
        input_path,
    ]

    try:
        subprocess.run(command, check=True)
    except FileNotFoundError as exc:
        raise SystemExit("Ghostscript is not installed or not reachable") from exc
    except subprocess.CalledProcessError as exc:
        raise SystemExit(f"Ghostscript failed with exit code {exc.returncode}") from exc

    return 0


def watermark_pdf(output_path: str, input_path: str, watermark_json: str) -> int:
    payload = json.loads(watermark_json)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    reader = PdfReader(input_path)
    writer = PdfWriter()

    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)

        with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
            overlay_path = Path(temp_directory) / "overlay.pdf"
            create_watermark_overlay(overlay_path, width, height, payload)
            overlay_reader = PdfReader(str(overlay_path))
            page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def add_page_numbers_pdf(output_path: str, input_path: str, options_json: str) -> int:
    payload = json.loads(options_json)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    reader = PdfReader(input_path)
    writer = PdfWriter()
    page_number_start = int(payload.get("pageNumberStart") or 1)

    for index, page in enumerate(reader.pages):
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
            overlay_path = Path(temp_directory) / "page-number-overlay.pdf"
            create_page_number_overlay(
                overlay_path,
                width,
                height,
                payload,
                page_number_start + index,
            )
            overlay_reader = PdfReader(str(overlay_path))
            page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def sign_stamp_pdf(output_path: str, input_path: str, options_json: str) -> int:
    payload = json.loads(options_json)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    reader = PdfReader(input_path)
    writer = PdfWriter()

    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
            overlay_path = Path(temp_directory) / "sign-overlay.pdf"
            create_sign_stamp_overlay(overlay_path, width, height, payload)
            overlay_reader = PdfReader(str(overlay_path))
            page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def batch_sign_stamp_pdf(zip_path: str, options_json: str, input_paths: list[str]) -> int:
    if len(input_paths) < 2:
        raise SystemExit("batch_sign_stamp_pdf requires at least two PDF files")

    zip_output = Path(zip_path)
    zip_output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(dir=str(zip_output.parent)) as temp_directory:
        temp_root = Path(temp_directory)
        generated_paths: list[Path] = []
        for input_path in input_paths:
            source = Path(input_path)
            output_path = temp_root / f"{source.stem}-stamped.pdf"
            sign_stamp_pdf(str(output_path), str(source), options_json)
            generated_paths.append(output_path)

        with zipfile.ZipFile(zip_output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for generated_path in generated_paths:
                archive.write(generated_path, arcname=generated_path.name)

    return 0


def rotate_pdf(output_path: str, input_path: str, rotation_angle: str) -> int:
    angle = int(rotation_angle)
    if angle not in {90, 180, 270}:
        raise SystemExit("rotate_pdf only supports 90, 180, or 270 degrees")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        rotated_page = page.rotate(angle)
        writer.add_page(rotated_page)

    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def create_watermark_overlay(output_path: Path, width: float, height: float, payload: dict) -> None:
    overlay = canvas.Canvas(str(output_path), pagesize=(width, height))
    if payload.get("watermarkType") == "image":
        draw_image_watermark(overlay, width, height, payload)
    else:
        draw_text_watermark(overlay, width, height, payload)
    overlay.save()


def create_page_number_overlay(
    output_path: Path, width: float, height: float, payload: dict, page_number: int
) -> None:
    overlay = canvas.Canvas(str(output_path), pagesize=(width, height))
    position = payload.get("pageNumberPosition") or "footer_center"
    page_number_format = payload.get("pageNumberFormat") or "cn_page"
    text = str(page_number) if page_number_format == "plain" else f"第 {page_number} 页"

    overlay.setFont("Helvetica", 11)
    y = 20
    if position == "bottom_right":
        overlay.drawRightString(width - 24, y, text)
    else:
        overlay.drawCentredString(width / 2, y, text)
    overlay.save()


def create_sign_stamp_overlay(output_path: Path, width: float, height: float, payload: dict) -> None:
    overlay = canvas.Canvas(str(output_path), pagesize=(width, height))
    draw_fixed_image_overlay(overlay, width, height, payload, payload.get("stampImagePath"))
    overlay.save()


def draw_text_watermark(overlay: canvas.Canvas, width: float, height: float, payload: dict) -> None:
    text_content = payload.get("textContent") or "水印"
    font_size = int(payload.get("fontSize") or 24)
    opacity = float(payload.get("opacity") or 0.18)
    rotation = float(payload.get("rotation") or -30)
    text_layout = payload.get("textLayout") or "tile"

    overlay.saveState()
    overlay.setFillAlpha(opacity)
    font_name = "STSong-Light" if any(ord(char) > 127 for char in text_content) else "Helvetica-Bold"
    overlay.setFont(font_name, font_size)

    if text_layout == "center":
        overlay.translate(width / 2, height / 2)
        overlay.rotate(rotation)
        overlay.drawCentredString(0, 0, text_content)
        overlay.restoreState()
        return

    step_x = max(font_size * 6, 160)
    step_y = max(font_size * 4, 120)
    overlay.translate(width / 2, height / 2)
    overlay.rotate(rotation)
    for x in range(-int(width), int(width) + step_x, step_x):
        for y in range(-int(height), int(height) + step_y, step_y):
            overlay.drawCentredString(x, y, text_content)
    overlay.restoreState()


def draw_image_watermark(overlay: canvas.Canvas, width: float, height: float, payload: dict) -> None:
    image_path = payload.get("watermarkImagePath")
    if not image_path:
        raise SystemExit("Image watermark requires watermarkImagePath")

    draw_fixed_image_overlay(overlay, width, height, payload, image_path)


def draw_fixed_image_overlay(
    overlay: canvas.Canvas, width: float, height: float, payload: dict, image_path: str
) -> None:
    image_scale_percent = max(5, min(100, int(payload.get("imageScalePercent") or 30)))
    opacity = float(payload.get("opacity") or 0.3)
    image_position = payload.get("imagePosition") or "center"

    with Image.open(image_path).convert("RGBA") as image:
        alpha = image.getchannel("A")
        alpha = alpha.point(lambda value: int(value * opacity))
        image.putalpha(alpha)

        max_target_width = width * (image_scale_percent / 100.0)
        scale = max_target_width / image.width
        target_width = max(1, int(image.width * scale))
        target_height = max(1, int(image.height * scale))
        resized = image.resize((target_width, target_height))

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_png:
            temp_png_path = Path(temp_png.name)
            resized.save(temp_png_path, format="PNG")

        try:
            margin = 24
            x = (width - target_width) / 2
            y = (height - target_height) / 2
            if image_position == "bottom_left":
                x = margin
                y = margin
            elif image_position == "bottom_right":
                x = width - target_width - margin
                y = margin

            overlay.drawImage(str(temp_png_path), x, y, width=target_width, height=target_height, mask='auto')
        finally:
            if temp_png_path.exists():
                temp_png_path.unlink()


def pdf_extract_pages(output_path: str, input_path: str, selection_json: str) -> int:
    payload = json.loads(selection_json)
    ordered_pages = payload.get("orderedPages") or []
    if not ordered_pages:
        raise SystemExit("pdf_extract_pages requires at least one selected page")

    reader = PdfReader(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)

    for page_number in ordered_pages:
        ensure_page_in_bounds(page_number, total_pages)
        writer.add_page(reader.pages[page_number - 1])

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("wb") as output_file:
        writer.write(output_file)

    return 0


def split_pdf(zip_path: str, input_path: str, selection_json: str) -> int:
    payload = json.loads(selection_json)
    outputs = payload.get("outputs") or []
    if not outputs:
        raise SystemExit("split_pdf requires at least one output range")

    reader = PdfReader(input_path)
    total_pages = len(reader.pages)
    source_name = Path(input_path).stem
    zip_output = Path(zip_path)
    zip_output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(dir=str(zip_output.parent)) as temp_directory:
        temp_root = Path(temp_directory)
        generated_paths: list[Path] = []

        for index, ordered_pages in enumerate(outputs, start=1):
            if not ordered_pages:
                raise SystemExit("split_pdf cannot create an empty output file")

            writer = PdfWriter()
            for page_number in ordered_pages:
                ensure_page_in_bounds(page_number, total_pages)
                writer.add_page(reader.pages[page_number - 1])

            part_path = temp_root / f"{source_name}-part-{index}.pdf"
            with part_path.open("wb") as part_file:
                writer.write(part_file)
            generated_paths.append(part_path)

        with zipfile.ZipFile(zip_output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for generated_path in generated_paths:
                archive.write(generated_path, arcname=generated_path.name)

    return 0


def zip_files(zip_path: str, input_paths: list[str]) -> int:
    if not input_paths:
        raise SystemExit("zip_files requires at least one input file")

    output = Path(zip_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for input_path in input_paths:
            path = Path(input_path)
            archive.write(path, arcname=path.name)

    return 0


def image_compress_batch(zip_path: str, options_json: str, input_paths: list[str]) -> int:
    if not input_paths:
        raise SystemExit("image_compress_batch requires at least one image")

    options = json.loads(options_json or "{}")
    quality = max(20, min(95, int(options.get("quality") or 75)))
    zip_output = Path(zip_path)
    zip_output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(dir=str(zip_output.parent)) as temp_directory:
        temp_root = Path(temp_directory)
        generated_paths: list[Path] = []
        for input_path in input_paths:
            source = Path(input_path)
            output_path = temp_root / f"{source.stem}-compressed{normalize_image_suffix(source.suffix, '.jpg')}"
            with Image.open(source) as image:
                save_image_file(image, output_path, quality=quality)
            generated_paths.append(output_path)

        write_files_to_zip(zip_output, generated_paths)

    return 0


def image_resize_exact(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    width = max(1, int(options.get("targetWidth") or 1))
    height = max(1, int(options.get("targetHeight") or 1))

    with Image.open(input_path) as image:
        resized = image.convert("RGBA").resize((width, height), Image.Resampling.LANCZOS)
        save_image_file(resized, Path(output_path))

    return 0


def image_resize_scale(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    scale_percent = max(1, int(options.get("scalePercent") or 100))

    with Image.open(input_path) as image:
        width = max(1, round(image.width * scale_percent / 100))
        height = max(1, round(image.height * scale_percent / 100))
        resized = image.convert("RGBA").resize((width, height), Image.Resampling.LANCZOS)
        save_image_file(resized, Path(output_path))

    return 0


def image_crop_free(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    crop_x = max(0, int(options.get("cropX") or 0))
    crop_y = max(0, int(options.get("cropY") or 0))
    crop_width = max(1, int(options.get("cropWidth") or 1))
    crop_height = max(1, int(options.get("cropHeight") or 1))

    with Image.open(input_path) as image:
        cropped = image.convert("RGBA").crop((
            crop_x,
            crop_y,
            min(image.width, crop_x + crop_width),
            min(image.height, crop_y + crop_height),
        ))
        save_image_file(cropped, Path(output_path))

    return 0


def image_crop_ratio(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    aspect_ratio = parse_ratio_text(options.get("aspectRatio") or "1:1")

    with Image.open(input_path) as image:
        cropped = center_crop_to_ratio(image.convert("RGBA"), aspect_ratio)
        save_image_file(cropped, Path(output_path))

    return 0


def image_crop_ratio_batch(zip_path: str, options_json: str, input_paths: list[str]) -> int:
    options = json.loads(options_json or "{}")
    aspect_ratio = parse_ratio_text(options.get("aspectRatio") or "1:1")
    output_format = normalize_format_name(options.get("outputFormat") or "png")
    zip_output = Path(zip_path)
    zip_output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(dir=str(zip_output.parent)) as temp_directory:
        temp_root = Path(temp_directory)
        generated_paths: list[Path] = []
        for input_path in input_paths:
            source = Path(input_path)
            output_ext = format_to_extension(output_format)
            output_path = temp_root / f"{source.stem}-ratio-cropped.{output_ext}"
            with Image.open(source) as image:
                cropped = center_crop_to_ratio(image.convert("RGBA"), aspect_ratio)
                save_image_file(cropped, output_path, force_format=output_format)
            generated_paths.append(output_path)

        write_files_to_zip(zip_output, generated_paths)

    return 0


def image_split_grid(zip_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    rows = max(1, int(options.get("rows") or 2))
    columns = max(1, int(options.get("columns") or 2))
    output_format = normalize_format_name(options.get("outputFormat") or "png")
    zip_output = Path(zip_path)
    zip_output.parent.mkdir(parents=True, exist_ok=True)

    with Image.open(input_path) as image:
        base_image = image.convert("RGBA")
        tile_width = max(1, math.floor(base_image.width / columns))
        tile_height = max(1, math.floor(base_image.height / rows))

        with tempfile.TemporaryDirectory(dir=str(zip_output.parent)) as temp_directory:
            temp_root = Path(temp_directory)
            generated_paths: list[Path] = []

            for row in range(rows):
                for column in range(columns):
                    left = column * tile_width
                    top = row * tile_height
                    right = base_image.width if column == columns - 1 else left + tile_width
                    bottom = base_image.height if row == rows - 1 else top + tile_height
                    part = base_image.crop((left, top, right, bottom))
                    output_path = temp_root / f"{Path(input_path).stem}-r{row + 1}-c{column + 1}.{format_to_extension(output_format)}"
                    save_image_file(part, output_path, force_format=output_format)
                    generated_paths.append(output_path)

            write_files_to_zip(zip_output, generated_paths)

    return 0


def image_concat_long(output_path: str, options_json: str, input_paths: list[str]) -> int:
    if not input_paths:
        raise SystemExit("image_concat_long requires at least one image")

    options = json.loads(options_json or "{}")
    direction = "horizontal" if options.get("direction") == "horizontal" else "vertical"
    gap = max(0, int(options.get("gap") or 0))
    background_color = options.get("backgroundColor") or "#ffffff"
    images = [Image.open(path).convert("RGBA") for path in input_paths]

    try:
        if direction == "horizontal":
            total_width = sum(image.width for image in images) + gap * max(0, len(images) - 1)
            total_height = max(image.height for image in images)
            canvas_image = Image.new("RGBA", (total_width, total_height), parse_color_value(background_color))
            cursor_x = 0
            for image in images:
                canvas_image.alpha_composite(image, (cursor_x, (total_height - image.height) // 2))
                cursor_x += image.width + gap
        else:
            total_width = max(image.width for image in images)
            total_height = sum(image.height for image in images) + gap * max(0, len(images) - 1)
            canvas_image = Image.new("RGBA", (total_width, total_height), parse_color_value(background_color))
            cursor_y = 0
            for image in images:
                canvas_image.alpha_composite(image, ((total_width - image.width) // 2, cursor_y))
                cursor_y += image.height + gap

        save_image_file(canvas_image, Path(output_path))
    finally:
        for image in images:
            image.close()

    return 0


def image_collage(output_path: str, options_json: str, input_paths: list[str]) -> int:
    if not input_paths:
        raise SystemExit("image_collage requires at least one image")

    options = json.loads(options_json or "{}")
    columns = max(1, int(options.get("columns") or 2))
    gap = max(0, int(options.get("gap") or 0))
    background_color = options.get("backgroundColor") or "#ffffff"
    images = [Image.open(path).convert("RGBA") for path in input_paths]

    try:
        cell_width = max(image.width for image in images)
        cell_height = max(image.height for image in images)
        rows = math.ceil(len(images) / columns)
        width = cell_width * columns + gap * max(0, columns - 1)
        height = cell_height * rows + gap * max(0, rows - 1)
        canvas_image = Image.new("RGBA", (width, height), parse_color_value(background_color))

        for index, image in enumerate(images):
            row = index // columns
            column = index % columns
            left = column * (cell_width + gap) + (cell_width - image.width) // 2
            top = row * (cell_height + gap) + (cell_height - image.height) // 2
            canvas_image.alpha_composite(image, (left, top))

        save_image_file(canvas_image, Path(output_path))
    finally:
        for image in images:
            image.close()

    return 0


def image_fill_background(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    background_color = parse_color_value(options.get("backgroundColor") or "#ffffff")

    with Image.open(input_path) as image:
        source = image.convert("RGBA")
        background = Image.new("RGBA", source.size, background_color)
        background.alpha_composite(source)
        save_image_file(background, Path(output_path), force_format=normalize_format_name(options.get("outputFormat") or Path(output_path).suffix))

    return 0


def image_watermark_tile(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    text_content = str(options.get("textContent") or "仅供内部使用")
    font_size = max(12, int(options.get("fontSize") or 24))
    opacity = max(0.05, min(0.9, float(options.get("opacity") or 0.22)))
    rotation = int(options.get("rotation") or -28)
    gap = max(40, int(options.get("gap") or 120))

    with Image.open(input_path) as image:
      base_image = image.convert("RGBA")
      overlay = Image.new("RGBA", base_image.size, (255, 255, 255, 0))
      drawer = ImageDraw.Draw(overlay)
      text_fill = (90, 90, 90, int(255 * opacity))
      step_x = max(gap, font_size * 4)
      step_y = max(gap, font_size * 3)
      for x in range(-base_image.width, base_image.width * 2, step_x):
          for y in range(-base_image.height, base_image.height * 2, step_y):
              drawer.text((x, y), text_content, fill=text_fill)

      rotated = overlay.rotate(rotation, expand=False)
      base_image.alpha_composite(rotated)
      save_image_file(base_image, Path(output_path))

    return 0


def image_grayscale(output_path: str, input_path: str, options_json: str) -> int:
    with Image.open(input_path) as image:
        grayscale = ImageOps.grayscale(image.convert("RGBA"))
        save_image_file(grayscale, Path(output_path), force_format=normalize_format_name(json.loads(options_json or "{}").get("outputFormat") or Path(output_path).suffix))
    return 0


def image_invert(output_path: str, input_path: str, options_json: str) -> int:
    with Image.open(input_path) as image:
        rgba = image.convert("RGBA")
        red, green, blue, alpha = rgba.split()
        rgb = Image.merge("RGB", (red, green, blue))
        inverted = ImageOps.invert(rgb)
        result = Image.merge("RGBA", (*inverted.split(), alpha))
        save_image_file(result, Path(output_path))
    return 0


def image_printmaking(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    threshold = max(0, min(255, int(options.get("threshold") or 126)))
    with Image.open(input_path) as image:
        grayscale = ImageOps.grayscale(image.convert("RGBA"))
        result = grayscale.point(lambda value: 255 if value > threshold else 0, mode="1").convert("L")
        save_image_file(result, Path(output_path), force_format=normalize_format_name(options.get("outputFormat") or Path(output_path).suffix))
    return 0


def image_emboss(output_path: str, input_path: str, options_json: str) -> int:
    with Image.open(input_path) as image:
        result = image.convert("RGBA").filter(ImageFilter.EMBOSS)
        save_image_file(result, Path(output_path))
    return 0


def image_remove_solid_bg(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    tolerance = max(0, min(255, int(options.get("tolerance") or 36)))

    with Image.open(input_path) as image:
        rgba = image.convert("RGBA")
        base_color = rgba.getpixel((0, 0))
        pixels = []
        for pixel in rgba.getdata():
            if all(abs(pixel[index] - base_color[index]) <= tolerance for index in range(3)):
                pixels.append((pixel[0], pixel[1], pixel[2], 0))
            else:
                pixels.append(pixel)
        rgba.putdata(pixels)
        save_image_file(rgba, Path(output_path), force_format="PNG")

    return 0


def image_smart_bg_remove(output_path: str, input_path: str, options_json: str) -> int:
    try:
        from rembg import new_session, remove
    except ImportError as exc:
        raise SystemExit("image_smart_bg_remove requires python package rembg[cpu]") from exc

    options = json.loads(options_json or "{}")
    model_name = str(
        options.get("modelName")
        or os.environ.get("REMBG_MODEL")
        or "isnet-general-use"
    ).strip() or "isnet-general-use"

    with open(input_path, "rb") as source_file:
        source_bytes = source_file.read()

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    result_bytes = remove(source_bytes, session=new_session(model_name))
    with Image.open(BytesIO(result_bytes)) as image:
        normalized = image.convert("RGBA")
        pixels = []
        for red, green, blue, alpha in normalized.getdata():
            next_alpha = 255 if alpha >= 250 else alpha
            pixels.append((red, green, blue, next_alpha))
        normalized.putdata(pixels)
        save_image_file(normalized, output, force_format="PNG")
    return 0


def qr_generate(output_path: str, options_json: str) -> int:
    import qrcode

    options = json.loads(options_json or "{}")
    qr_text = str(options.get("qrText") or "").strip()
    if not qr_text:
        raise SystemExit("qr_generate requires qrText")
    size_px = max(128, min(1024, int(options.get("sizePx") or 320)))

    qr = qrcode.QRCode(border=2, box_size=10)
    qr.add_data(qr_text)
    qr.make(fit=True)
    image = qr.make_image(fill_color="black", back_color="white").convert("RGBA")
    image = image.resize((size_px, size_px), Image.Resampling.NEAREST)
    save_image_file(image, Path(output_path), force_format="PNG")
    return 0


def qr_generate_batch(zip_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    lines = [line.strip() for line in str(options.get("qrLinesText") or "").splitlines() if line.strip()]
    if not lines:
        raise SystemExit("qr_generate_batch requires at least one line")

    zip_output = Path(zip_path)
    zip_output.parent.mkdir(parents=True, exist_ok=True)
    size_px = max(128, min(1024, int(options.get("sizePx") or 256)))

    with tempfile.TemporaryDirectory(dir=str(zip_output.parent)) as temp_directory:
        temp_root = Path(temp_directory)
        generated_paths: list[Path] = []
        for index, line in enumerate(lines, start=1):
            output_path = temp_root / f"qr-{index:03d}.png"
            qr_generate(str(output_path), json.dumps({"qrText": line, "sizePx": size_px}))
            generated_paths.append(output_path)

        write_files_to_zip(zip_output, generated_paths)

    return 0


def qr_decode(output_path: str, input_path: str) -> int:
    detector = cv2.QRCodeDetector()
    image = cv2.imread(input_path)
    if image is None:
        raise SystemExit("qr_decode could not read the input image")

    decoded_text = ""
    try:
        decoded_text, points, _ = detector.detectAndDecode(image)
        if not decoded_text:
            success, decoded_info, _, _ = detector.detectAndDecodeMulti(image)
            if success:
                decoded_text = "\n".join([item for item in decoded_info if item])
    except Exception as exc:
        raise SystemExit(f"qr_decode failed: {exc}") from exc

    if not decoded_text.strip():
        raise SystemExit("未识别到二维码内容")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(decoded_text.strip(), encoding="utf8")
    return 0


def image_heic_convert(output_path: str, input_path: str, options_json: str) -> int:
    try:
        from pillow_heif import register_heif_opener
    except ImportError as exc:
        raise SystemExit("image_heic_convert requires python package pillow-heif") from exc

    options = json.loads(options_json or "{}")
    output_format = normalize_format_name(options.get("outputFormat") or "jpg")
    if output_format not in {"JPEG", "PNG"}:
        output_format = "JPEG"

    register_heif_opener()
    with Image.open(input_path) as image:
        source = ImageOps.exif_transpose(image)
        save_image_file(source, Path(output_path), force_format=output_format)

    return 0


def favicon_generate(output_path: str, input_path: str, _options_json: str) -> int:
    with Image.open(input_path) as image:
        icon = image.convert("RGBA")
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        icon.save(output, format="ICO", sizes=[(16, 16), (32, 32), (48, 48), (64, 64)])
    return 0


def app_icon_generate(output_path: str, input_path: str, _options_json: str) -> int:
    return export_icon_pack(output_path, input_path, [16, 32, 48, 64, 128, 256, 512], "app-icon")


def chrome_icon_generate(output_path: str, input_path: str, _options_json: str) -> int:
    return export_icon_pack(output_path, input_path, [16, 32, 48, 64, 128, 256], "chrome-icon")


def image_add_padding(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    top = max(0, int(options.get("paddingTop") or 0))
    right = max(0, int(options.get("paddingRight") or 0))
    bottom = max(0, int(options.get("paddingBottom") or 0))
    left = max(0, int(options.get("paddingLeft") or 0))
    background_color = parse_color_value(options.get("backgroundColor") or "#ffffff")

    with Image.open(input_path) as image:
        source = image.convert("RGBA")
        target = Image.new("RGBA", (source.width + left + right, source.height + top + bottom), background_color)
        target.alpha_composite(source, (left, top))
        save_image_file(target, Path(output_path))

    return 0


def image_pixelate(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    block_size = max(2, int(options.get("blockSize") or 12))
    with Image.open(input_path) as image:
        source = image.convert("RGBA")
        reduced = source.resize((max(1, source.width // block_size), max(1, source.height // block_size)), Image.Resampling.BILINEAR)
        result = reduced.resize(source.size, Image.Resampling.NEAREST)
        save_image_file(result, Path(output_path))
    return 0


def image_increase_size(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    target_size_kb = max(1, int(options.get("targetSizeKb") or 100))
    target_size_bytes = target_size_kb * 1024
    with Image.open(input_path) as image:
        save_image_file(image.convert("RGBA"), Path(output_path))

    output = Path(output_path)
    current_size = output.stat().st_size
    if current_size < target_size_bytes:
        with output.open("ab") as output_file:
            output_file.write(b"\0" * (target_size_bytes - current_size))
    return 0


def image_clear_content(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    transparent = bool(options.get("transparent"))
    background = (0, 0, 0, 0) if transparent else parse_color_value(options.get("backgroundColor") or "#ffffff")
    with Image.open(input_path) as image:
        cleared = Image.new("RGBA", image.size, background)
        save_image_file(cleared, Path(output_path), force_format="PNG" if transparent else normalize_format_name(options.get("outputFormat") or Path(output_path).suffix))
    return 0


def image_format_convert_dispatch(output_path: str, args: list[str]) -> int:
    if not args:
        raise SystemExit("image_format_convert requires arguments")

    if args[0].strip().startswith("{"):
        return image_format_convert(output_path, args[0], args[1:])

    options_json = args[1] if len(args) > 1 else "{}"
    return image_format_convert(output_path, options_json, [args[0]])


def image_format_convert(output_path: str, options_json: str, input_paths: list[str]) -> int:
    if not input_paths:
        raise SystemExit("image_format_convert requires at least one image")

    options = json.loads(options_json or "{}")
    output_format = normalize_format_name(options.get("outputFormat") or "png")
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    if output.suffix.lower() != ".zip" and len(input_paths) == 1:
        with Image.open(input_paths[0]) as image:
            save_image_file(image.convert("RGBA"), output, force_format=output_format)
        return 0

    with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
        temp_root = Path(temp_directory)
        generated_paths: list[Path] = []
        for input_path in input_paths:
            source = Path(input_path)
            converted_path = temp_root / f"{source.stem}-converted.{format_to_extension(output_format)}"
            with Image.open(source) as image:
                save_image_file(image.convert("RGBA"), converted_path, force_format=output_format)
            generated_paths.append(converted_path)

        write_files_to_zip(output, generated_paths)
    return 0


def extract_images_from_office_zip(output_path: str, input_path: str, media_prefix: str) -> int:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
        temp_root = Path(temp_directory)
        extracted_paths: list[Path] = []
        with zipfile.ZipFile(input_path) as archive:
            for name in archive.namelist():
                if not name.startswith(media_prefix):
                    continue
                content = archive.read(name)
                target_path = temp_root / Path(name).name
                target_path.write_bytes(content)
                extracted_paths.append(target_path)

        if not extracted_paths:
            raise SystemExit("No images found in the uploaded Office file")

        write_files_to_zip(output, extracted_paths)
    return 0


def image_modify_dpi(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    dpi = max(1, int(options.get("dpi") or 300))
    output_format = normalize_format_name(options.get("outputFormat") or Path(output_path).suffix)
    with Image.open(input_path) as image:
        save_image_file(image.convert("RGBA"), Path(output_path), force_format=output_format, dpi=(dpi, dpi))
    return 0


def gif_split(output_path: str, input_path: str, _options_json: str) -> int:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    with Image.open(input_path) as image:
        if getattr(image, "n_frames", 1) <= 1:
            raise SystemExit("GIF only has one frame")

        with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
            temp_root = Path(temp_directory)
            generated_paths: list[Path] = []
            for frame_index in range(image.n_frames):
                image.seek(frame_index)
                frame = image.convert("RGBA")
                frame_path = temp_root / f"{Path(input_path).stem}-frame-{frame_index + 1}.png"
                frame.save(frame_path, format="PNG")
                generated_paths.append(frame_path)

            write_files_to_zip(output, generated_paths)

    return 0


def gif_merge(output_path: str, options_json: str, input_paths: list[str]) -> int:
    if not input_paths:
        raise SystemExit("gif_merge requires at least one image")

    options = json.loads(options_json or "{}")
    duration_ms = max(50, int(options.get("durationMs") or 400))
    frames = [Image.open(path).convert("RGBA") for path in input_paths]

    try:
        first_frame = frames[0]
        other_frames = frames[1:]
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        first_frame.save(
            output,
            format="GIF",
            save_all=True,
            append_images=other_frames,
            duration=duration_ms,
            loop=0,
            disposal=2,
        )
    finally:
        for frame in frames:
            frame.close()
    return 0


def png_alpha_invert(output_path: str, input_path: str, _options_json: str) -> int:
    with Image.open(input_path) as image:
        rgba = image.convert("RGBA")
        red, green, blue, alpha = rgba.split()
        inverted_alpha = ImageOps.invert(alpha)
        result = Image.merge("RGBA", (red, green, blue, inverted_alpha))
        save_image_file(result, Path(output_path), force_format="PNG")
    return 0


def image_round_corner(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    radius = max(2, int(options.get("radius") or 36))
    with Image.open(input_path) as image:
        source = image.convert("RGBA")
        mask = Image.new("L", source.size, 0)
        drawer = ImageDraw.Draw(mask)
        drawer.rounded_rectangle((0, 0, source.width, source.height), radius=radius, fill=255)
        source.putalpha(mask)
        save_image_file(source, Path(output_path), force_format="PNG")
    return 0


def image_tile_fill(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    target_width = max(1, int(options.get("targetWidth") or 1200))
    target_height = max(1, int(options.get("targetHeight") or 1200))
    output_format = normalize_format_name(options.get("outputFormat") or Path(output_path).suffix)
    with Image.open(input_path) as image:
        tile = image.convert("RGBA")
        target = Image.new("RGBA", (target_width, target_height), (255, 255, 255, 0))
        for x in range(0, target_width, tile.width):
            for y in range(0, target_height, tile.height):
                target.alpha_composite(tile, (x, y))
        save_image_file(target, Path(output_path), force_format=output_format)
    return 0


def id_photo_resize(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    preset = resolve_id_photo_preset(options.get("idPhotoPreset") or "one_inch")
    max_size_kb = max(20, int(options.get("maxSizeKb") or 120))
    output_format = normalize_format_name(options.get("outputFormat") or "jpg")

    with Image.open(input_path) as image:
        cropped = center_crop_to_ratio(image.convert("RGBA"), preset[0] / preset[1])
        resized = cropped.resize(preset, Image.Resampling.LANCZOS)
        save_image_file(resized, Path(output_path), force_format=output_format, max_size_kb=max_size_kb)
    return 0


def id_photo_crop(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    preset = resolve_id_photo_preset(options.get("idPhotoPreset") or "one_inch")
    output_format = normalize_format_name(options.get("outputFormat") or "jpg")
    with Image.open(input_path) as image:
        cropped = center_crop_to_ratio(image.convert("RGBA"), preset[0] / preset[1])
        resized = cropped.resize(preset, Image.Resampling.LANCZOS)
        save_image_file(resized, Path(output_path), force_format=output_format)
    return 0


def id_photo_bg_swap(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    background_color = parse_color_value(options.get("backgroundColor") or "#438edb")
    tolerance = max(0, min(255, int(options.get("tolerance") or 36)))
    output_format = normalize_format_name(options.get("outputFormat") or "jpg")

    with Image.open(input_path) as image:
        rgba = image.convert("RGBA")
        base_color = rgba.getpixel((0, 0))
        pixels = []
        for pixel in rgba.getdata():
            if all(abs(pixel[index] - base_color[index]) <= tolerance for index in range(3)):
                pixels.append((*background_color[:3], 255))
            else:
                pixels.append(pixel)
        rgba.putdata(pixels)
        save_image_file(rgba, Path(output_path), force_format=output_format)
    return 0


def anti_ocr_image(output_path: str, input_path: str, options_json: str) -> int:
    options = json.loads(options_json or "{}")
    noise_level = max(1, min(80, int(options.get("noiseLevel") or 18)))
    output_format = normalize_format_name(options.get("outputFormat") or Path(output_path).suffix)

    with Image.open(input_path) as image:
        source = image.convert("RGBA")
        drawer = ImageDraw.Draw(source)
        for _ in range(max(12, (source.width * source.height) // 8000)):
            x = random.randint(0, max(0, source.width - 1))
            y = random.randint(0, max(0, source.height - 1))
            delta = random.randint(-noise_level, noise_level)
            red, green, blue, alpha = source.getpixel((x, y))
            source.putpixel((x, y), (
                clamp_color(red + delta),
                clamp_color(green - delta),
                clamp_color(blue + delta // 2),
                alpha,
            ))

        for _ in range(6):
            start = (random.randint(0, source.width), random.randint(0, source.height))
            end = (random.randint(0, source.width), random.randint(0, source.height))
            drawer.line((start, end), fill=(120, 120, 120, 72), width=1)

        save_image_file(source, Path(output_path), force_format=output_format)
    return 0


def export_icon_pack(output_path: str, input_path: str, sizes: list[int], name_prefix: str) -> int:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    with Image.open(input_path) as image:
        source = image.convert("RGBA")
        with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
            temp_root = Path(temp_directory)
            generated_paths: list[Path] = []
            for size in sizes:
                resized = source.resize((size, size), Image.Resampling.LANCZOS)
                icon_path = temp_root / f"{name_prefix}-{size}.png"
                resized.save(icon_path, format="PNG")
                generated_paths.append(icon_path)
            write_files_to_zip(output, generated_paths)
    return 0


def write_files_to_zip(zip_path: Path, input_paths: list[Path]) -> None:
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for input_path in input_paths:
            archive.write(input_path, arcname=input_path.name)


def save_image_file(
    image: Image.Image,
    output_path: Path,
    force_format: str | None = None,
    quality: int = 92,
    dpi: tuple[int, int] | None = None,
    max_size_kb: int | None = None,
) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    format_name = normalize_format_name(force_format or output_path.suffix)

    image_to_save = image.copy()
    save_options: dict = {}
    if dpi:
        save_options["dpi"] = dpi

    if format_name in {"JPEG", "JPG"}:
        image_to_save = image_to_save.convert("RGB")
        save_options.update({"format": "JPEG", "quality": quality, "optimize": True})
    elif format_name == "PNG":
        image_to_save = image_to_save.convert("RGBA")
        save_options.update({"format": "PNG", "optimize": True})
    elif format_name == "WEBP":
        save_options.update({"format": "WEBP", "quality": quality, "method": 6})
    elif format_name == "GIF":
        image_to_save = image_to_save.convert("P", palette=Image.ADAPTIVE)
        save_options.update({"format": "GIF"})
    else:
        save_options.update({"format": format_name})

    if max_size_kb and save_options.get("format") == "JPEG":
        current_quality = min(95, max(20, quality))
        while current_quality >= 20:
            image_to_save.save(output_path, **{**save_options, "quality": current_quality})
            if output_path.stat().st_size <= max_size_kb * 1024 or current_quality == 20:
                break
            current_quality -= 5
        return

    image_to_save.save(output_path, **save_options)


def normalize_format_name(value: str | None) -> str:
    normalized = str(value or "").strip().lower().lstrip(".")
    if normalized in {"jpg", "jpeg"}:
        return "JPEG"
    if normalized == "png":
        return "PNG"
    if normalized == "webp":
        return "WEBP"
    if normalized == "gif":
        return "GIF"
    if normalized == "ico":
        return "ICO"
    return "PNG"


def format_to_extension(format_name: str) -> str:
    normalized = normalize_format_name(format_name)
    if normalized == "JPEG":
        return "jpg"
    return normalized.lower()


def normalize_image_suffix(current_suffix: str, fallback_suffix: str) -> str:
    suffix = current_suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return ".jpg" if suffix == ".jpeg" else suffix
    return fallback_suffix


def parse_ratio_text(value: str) -> float:
    text = str(value or "1:1").strip()
    if ":" not in text:
        return 1.0
    left_text, right_text = text.split(":", 1)
    left = max(1.0, float(left_text or 1))
    right = max(1.0, float(right_text or 1))
    return left / right


def center_crop_to_ratio(image: Image.Image, aspect_ratio: float) -> Image.Image:
    source_ratio = image.width / max(1, image.height)
    if source_ratio > aspect_ratio:
        target_width = int(image.height * aspect_ratio)
        left = max(0, (image.width - target_width) // 2)
        return image.crop((left, 0, left + target_width, image.height))
    target_height = int(image.width / max(0.0001, aspect_ratio))
    top = max(0, (image.height - target_height) // 2)
    return image.crop((0, top, image.width, top + target_height))


def parse_color_value(value: str) -> tuple[int, int, int, int]:
    color_text = str(value or "#ffffff").strip()
    if not color_text.startswith("#") or len(color_text) not in {4, 7}:
        color_text = "#ffffff"
    if len(color_text) == 4:
        color_text = "#" + "".join(char * 2 for char in color_text[1:])
    red = int(color_text[1:3], 16)
    green = int(color_text[3:5], 16)
    blue = int(color_text[5:7], 16)
    return red, green, blue, 255


def resolve_id_photo_preset(preset_key: str) -> tuple[int, int]:
    presets = {
        "one_inch": (295, 413),
        "two_inch": (413, 579),
        "small_one_inch": (260, 378),
    }
    return presets.get(str(preset_key or "one_inch"), presets["one_inch"])


def clamp_color(value: int) -> int:
    return max(0, min(255, int(value)))


def estimate_edge_background_color(image: Image.Image) -> tuple[int, int, int]:
    edge_pixels: list[tuple[int, int, int]] = []
    width, height = image.size
    pixels = image.load()
    for x in range(width):
        edge_pixels.append(pixels[x, 0][:3])
        edge_pixels.append(pixels[x, height - 1][:3])
    for y in range(height):
        edge_pixels.append(pixels[0, y][:3])
        edge_pixels.append(pixels[width - 1, y][:3])

    red = round(sum(pixel[0] for pixel in edge_pixels) / max(1, len(edge_pixels)))
    green = round(sum(pixel[1] for pixel in edge_pixels) / max(1, len(edge_pixels)))
    blue = round(sum(pixel[2] for pixel in edge_pixels) / max(1, len(edge_pixels)))
    return red, green, blue


def pixel_matches_background(
    pixel: tuple[int, int, int, int],
    background_color: tuple[int, int, int],
    tolerance: int,
) -> bool:
    if pixel[3] <= 16:
        return True

    return (
        abs(pixel[0] - background_color[0]) <= tolerance
        and abs(pixel[1] - background_color[1]) <= tolerance
        and abs(pixel[2] - background_color[2]) <= tolerance
    )


def sanitize_archive_stem(value: str) -> str:
    return (
        str(value or "file")
        .replace("/", "_")
        .replace("\\", "_")
        .replace(":", "_")
        .replace("*", "_")
        .replace("?", "_")
        .replace('"', "_")
        .replace("<", "_")
        .replace(">", "_")
        .replace("|", "_")
        .strip()
        or "file"
    )


def ensure_unique_archive_name(file_name: str, used_names: set[str]) -> str:
    candidate = file_name
    stem = Path(file_name).stem
    suffix = Path(file_name).suffix
    serial = 2
    while candidate in used_names:
        candidate = f"{stem}-{serial}{suffix}"
        serial += 1
    used_names.add(candidate)
    return candidate


def text_to_speech(
    output_path: str,
    source_text: str,
    language: str,
    output_format: str,
    ffmpeg_bin: str,
) -> int:
    try:
        import edge_tts
    except ImportError as exc:
        raise SystemExit("text_to_speech requires python package edge-tts") from exc

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    voice_map = {
        "zh": "zh-CN-XiaoxiaoNeural",
        "en": "en-US-AriaNeural",
    }
    voice_name = voice_map.get(language, voice_map["zh"])

    async def synthesize_mp3(mp3_path: Path) -> None:
        communicator = edge_tts.Communicate(source_text, voice_name)
        await communicator.save(str(mp3_path))

    if output_format == "mp3":
        asyncio.run(synthesize_mp3(output))
        return 0

    if output_format == "wav":
        if not ffmpeg_bin:
            raise SystemExit("text_to_speech wav output requires ffmpeg")

        with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
            temp_mp3 = Path(temp_directory) / "tts-output.mp3"
            asyncio.run(synthesize_mp3(temp_mp3))
            try:
                subprocess.run(
                    [
                        ffmpeg_bin,
                        "-y",
                        "-i",
                        str(temp_mp3),
                        str(output),
                    ],
                    check=True,
                )
            except FileNotFoundError as exc:
                raise SystemExit("ffmpeg is not installed or not reachable") from exc
            except subprocess.CalledProcessError as exc:
                raise SystemExit(f"ffmpeg failed with exit code {exc.returncode}") from exc
        return 0

    raise SystemExit("text_to_speech only supports mp3 or wav output")


def audio_to_text(
    output_path: str,
    input_path: str,
    language: str,
    ffmpeg_bin: str,
) -> int:
    try:
        import numpy as np
        import whisper
    except ImportError as exc:
        raise SystemExit("audio_to_text requires python package openai-whisper") from exc

    if not ffmpeg_bin:
        raise SystemExit("audio_to_text requires ffmpeg")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(dir=str(output.parent)) as temp_directory:
        temp_wav = Path(temp_directory) / "audio-input.wav"
        try:
            subprocess.run(
                [
                    ffmpeg_bin,
                    "-y",
                    "-i",
                    input_path,
                    "-ac",
                    "1",
                    "-ar",
                    "16000",
                    str(temp_wav),
                ],
                check=True,
            )
        except FileNotFoundError as exc:
            raise SystemExit("ffmpeg is not installed or not reachable") from exc
        except subprocess.CalledProcessError as exc:
            raise SystemExit(f"ffmpeg failed with exit code {exc.returncode}") from exc

        with wave.open(str(temp_wav), "rb") as wave_file:
            frames = wave_file.readframes(wave_file.getnframes())
            audio_array = np.frombuffer(frames, dtype=np.int16).astype(np.float32) / 32768.0

        model = whisper.load_model(os.environ.get("WHISPER_MODEL", "base"))
        result = model.transcribe(
            audio_array,
            language=None if language == "auto" else language,
            fp16=False,
        )
        text = str(result.get("text") or "").strip()
        if not text:
            raise SystemExit("audio_to_text did not produce any text")

        output.write_text(text, encoding="utf8")

    return 0


def save_image_file(
    image: Image.Image,
    output_path: Path,
    force_format: str | None = None,
    quality: int = 92,
    dpi: tuple[int, int] | None = None,
    max_size_kb: int | None = None,
) -> None:
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    format_name = normalize_format_name(force_format or output.suffix or "png")

    save_kwargs: dict = {}
    if dpi:
        save_kwargs["dpi"] = dpi

    source_image = image.convert("RGBA") if image.mode not in {"RGB", "RGBA", "L"} else image.copy()
    image_to_save = source_image

    if format_name == "JPEG":
        background = Image.new("RGB", source_image.size, (255, 255, 255))
        if source_image.mode == "RGBA":
            background.paste(source_image, mask=source_image.getchannel("A"))
        else:
            background.paste(source_image.convert("RGB"))
        image_to_save = background
        save_kwargs["quality"] = max(20, min(95, int(quality)))
        save_kwargs["optimize"] = True
    elif format_name == "WEBP":
        image_to_save = source_image.convert("RGBA")
        save_kwargs["quality"] = max(20, min(95, int(quality)))
        save_kwargs["method"] = 6
    elif format_name == "PNG":
        image_to_save = source_image.convert("RGBA")
        save_kwargs["optimize"] = True
    else:
        image_to_save = source_image

    image_to_save.save(output, format=format_name, **save_kwargs)

    if max_size_kb and format_name in {"JPEG", "WEBP"}:
        target_size_bytes = max_size_kb * 1024
        current_quality = int(save_kwargs.get("quality", quality))
        while output.stat().st_size > target_size_bytes and current_quality > 20:
            current_quality -= 5
            image_to_save.save(output, format=format_name, quality=current_quality, **{
                key: value for key, value in save_kwargs.items() if key != "quality"
            })


def normalize_format_name(value: str) -> str:
    normalized = str(value or "").strip().lower().lstrip(".")
    mapping = {
        "jpg": "JPEG",
        "jpeg": "JPEG",
        "png": "PNG",
        "webp": "WEBP",
        "gif": "GIF",
        "bmp": "BMP",
        "ico": "ICO",
        "tif": "TIFF",
        "tiff": "TIFF",
    }
    return mapping.get(normalized, "PNG")


def format_to_extension(format_name: str) -> str:
    normalized = normalize_format_name(format_name)
    mapping = {
        "JPEG": "jpg",
        "PNG": "png",
        "WEBP": "webp",
        "GIF": "gif",
        "BMP": "bmp",
        "ICO": "ico",
        "TIFF": "tiff",
    }
    return mapping.get(normalized, "png")


def ensure_page_in_bounds(page_number: int, total_pages: int) -> None:
    if page_number < 1 or page_number > total_pages:
        raise SystemExit(f"第 {page_number} 页超出 PDF 总页数（共 {total_pages} 页）。")


def extract_docx_title(input_path: str) -> str:
    path = Path(input_path)
    try:
        with zipfile.ZipFile(path) as archive:
            with archive.open("docProps/core.xml") as core_file:
                tree = ElementTree.parse(core_file)
    except Exception:
        return ""

    namespace = {"dc": "http://purl.org/dc/elements/1.1/"}
    title_node = tree.find(".//dc:title", namespace)
    return title_node.text.strip() if title_node is not None and title_node.text else ""


def wrap_text(text: str, max_chars: int) -> list[str]:
    if len(text) <= max_chars:
        return [text]

    words = text.split()
    if not words:
        return [text]

    lines: list[str] = []
    current = ""
    for word in words:
        candidate = word if not current else current + " " + word
        if len(candidate) <= max_chars:
            current = candidate
            continue
        if current:
            lines.append(current)
        current = word
    if current:
        lines.append(current)
    return lines or [text]


if __name__ == "__main__":
    raise SystemExit(main())
